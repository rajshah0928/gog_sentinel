"""
Builds the Solution Presentation (PPTX) deliverable from docs/presentation_outline.md's
content, with real evidence screenshots embedded. Run once, from repo root:

    python docs/build_presentation.py

Produces docs/Sentinel_Solution_Presentation.pptx. Not part of the running
application - a one-off document-generation script, kept in docs/ alongside
its output rather than in the main package tree.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOCS_DIR = Path(__file__).resolve().parent
EVIDENCE_DIR = DOCS_DIR / "evidence"
OUT_PATH = DOCS_DIR / "Sentinel_Solution_Presentation.pptx"
LOGO_PATH = EVIDENCE_DIR / "ruvision_logo_small.png"
COMPANY_NAME = "RuVision Thinking Labs"

# --- palette (matches the dashboard's dark theme for visual consistency) ---
BG = RGBColor(0x0A, 0x0E, 0x14)
BG_CARD = RGBColor(0x11, 0x18, 0x23)
BORDER = RGBColor(0x1E, 0x2A, 0x3A)
TEXT_HI = RGBColor(0xE6, 0xED, 0xF5)
TEXT_MID = RGBColor(0x9F, 0xB0, 0xC6)
TEXT_LO = RGBColor(0x5C, 0x6B, 0x86)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide(logo=True):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    if logo and LOGO_PATH.exists():
        add_logo(slide)
    return slide


def add_logo(slide, height=Inches(0.4)):
    """
    Top-right RuVision brand mark, present on every slide - placement
    matches the company's existing deck convention (logo pinned top-right,
    ~0.2in from the top edge, ~0.4in tall).
    """
    from PIL import Image
    im = Image.open(LOGO_PATH)
    ratio = im.width / im.height
    width = Emu(int(height * ratio))
    left = SLIDE_W - width - Inches(0.3)
    top = Inches(0.2)
    slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width, height=height)


def add_text(slide, left, top, width, height, text, size=18, color=TEXT_HI,
             bold=False, align=PP_ALIGN.LEFT, font="Inter", anchor=None, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=TEXT_MID,
                 bold_color=TEXT_HI, font="Inter", space_after=10, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.level = level
        bullet = "▸ " if level == 0 else "–  "
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_header(slide, kicker, title):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35),
              kicker.upper(), size=12, color=ACCENT, bold=True, font="Consolas")
    add_text(slide, Inches(0.6), Inches(0.68), Inches(11), Inches(0.7),
              title, size=28, color=TEXT_HI, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(12.1), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    line.shadow.inherit = False


def add_card(slide, left, top, width, height, fill=BG_CARD, line_color=BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = line_color
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def add_image_framed(slide, path, left, top, width, caption=None):
    from PIL import Image
    im = Image.open(path)
    ratio = im.height / im.width
    height = Emu(int(width * ratio))
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    # thin border
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1.5)
    if caption:
        add_text(slide, left, top + height + Inches(0.06), width, Inches(0.3),
                  caption, size=11, color=TEXT_LO, align=PP_ALIGN.CENTER, font="Consolas")
    return pic, height


def add_footer(slide, page_num):
    add_text(slide, Inches(0.6), Inches(7.08), Inches(6), Inches(0.3),
              "SENTINEL — Gujarat Police Innovation Challenge 2026 · Model 2", size=9, color=TEXT_LO, font="Consolas")
    add_text(slide, Inches(12.0), Inches(7.08), Inches(0.7), Inches(0.3),
              str(page_num), size=9, color=TEXT_LO, align=PP_ALIGN.RIGHT, font="Consolas")


# ============================================================ Slide 1: Title
s = add_slide()
add_text(s, Inches(0), Inches(2.6), SLIDE_W, Inches(1.2),
          "SENTINEL", size=64, color=TEXT_HI, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.75), SLIDE_W, Inches(0.6),
          "Unified Viewing & Metadata Analytics", size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.35), SLIDE_W, Inches(0.5),
          "Gujarat Police Innovation Challenge 2026  ·  Model 2", size=16, color=TEXT_MID, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, Inches(0), Inches(5.1), SLIDE_W, Inches(0.4),
          f"Built by {COMPANY_NAME}", size=14, color=TEXT_LO, align=PP_ALIGN.CENTER, font="Consolas")
dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.55), Inches(2.35), Inches(0.16), Inches(0.16))
dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background(); dot.shadow.inherit = False

# ============================================================ Slide 2: Problem
s = add_slide()
add_header(s, "Background", "The Problem")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.5), Inches(4.5), [
    "26 Government Departments operate independent CCTV systems — different vendors, VMS platforms, storage architectures, and retention periods (7 to 15+ days).",
    "Central command centres must access these feeds through multiple separate viewer systems — no single pane of glass.",
    "High-value law-enforcement databases already exist and are digitized — VAHAN (vehicle registration), eGujCop / CCTNS (stolen vehicles, wanted persons), AFIS / NAFIS — but they are siloed from live video.",
    "Result: no automated, real-time correlation between what a camera sees and what these databases already know.",
], size=17, space_after=18)
add_footer(s, 2)

# ============================================================ Slide 3: Model choice
s = add_slide()
add_header(s, "Solution Model", "Model 2 — Unified Viewing & Metadata Analytics")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(7.0), Inches(4.8), [
    "Direct RTSP / ONVIF / vendor-API connection to each camera or VMS.",
    "No middleware or federation layer introduced.",
    "Departmental VMS and storage systems remain completely untouched and continue operating independently.",
], size=17, space_after=16)

card = add_card(s, Inches(8.0), Inches(1.6), Inches(4.7), Inches(4.8))
add_text(s, Inches(8.3), Inches(1.85), Inches(4.1), Inches(0.4), "WHY MODEL 2, NOT 3/4", size=13, color=ACCENT, bold=True, font="Consolas")
add_bullets(s, Inches(8.3), Inches(2.35), Inches(4.15), Inches(3.9), [
    "26 departments' VMS/storage diversity makes a federation layer a multi-department program — not a hackathon-timeline build.",
    "Delivers exactly the two capabilities the evaluation tests: live unified view + AI vehicle tracing.",
    "Lowest integration risk: zero dependency on any department changing its systems.",
], size=14, space_after=12)
add_footer(s, 3)

# ============================================================ Slide 4: Architecture
s = add_slide()
add_header(s, "Architecture", "End-to-End Data Flow")

stages = [
    ("Departmental\nCCTV / VMS", TEXT_MID),
    ("Capture Layer", ACCENT),
    ("ANPR\nAnalytics", ACCENT),
    ("Watchlist &\nAlerting", AMBER),
    ("Route\nReconstruction", ACCENT),
    ("Unified\nDashboard", GREEN),
]
box_w = Inches(1.85)
gap = Inches(0.28)
total_w = box_w * len(stages) + gap * (len(stages) - 1)
start_x = int((SLIDE_W - total_w) / 2)
y = Inches(2.6)
x = start_x
for i, (label, color) in enumerate(stages):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, Inches(1.15))
    card.adjustments[0] = 0.12
    card.fill.solid(); card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(12.5); run.font.bold = True; run.font.color.rgb = TEXT_HI
    if i < len(stages) - 1:
        arrow_x = x + box_w
        arrow = add_text(s, arrow_x, y + Inches(0.35), gap, Inches(0.5), "→", size=20, color=TEXT_LO, align=PP_ALIGN.CENTER)
    x = x + box_w + gap

add_text(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(0.5),
          "No video is ever centrally stored — only structured detection metadata (plate text, camera, timestamp, confidence).",
          size=15, color=TEXT_MID, align=PP_ALIGN.CENTER)

card2 = add_card(s, Inches(2.0), Inches(5.1), Inches(9.3), Inches(1.5))
add_bullets(s, Inches(2.3), Inches(5.3), Inches(8.8), Inches(1.2), [
    "Capture: TCP-forced RTSP, PTS-derived timing, exponential-backoff reconnect.",
    "Analytics: YOLO plate detector + PaddleOCR, CPU-only, configurable sample rate.",
    "Storage: SQLite metadata log (schema-compatible upgrade path to PostgreSQL).",
], size=13.5, space_after=8)
add_footer(s, 4)

# ============================================================ Slide 5: Live stream ingestion
s = add_slide()
add_header(s, "Engineering — Proven, Not Just Designed", "Live Stream Ingestion")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.9), Inches(2.0), [
    "TCP-forced RTSP, PTS-derived timing (never wall-clock or declared FPS), exponential backoff reconnect.",
    "Tolerant of decoder warm-up noise and mixed H.264 / H.265 codecs — confirmed live on the sandbox grid.",
    "Camera inventory always pulled live from the gateway catalogue — zero hardcoded camera IDs.",
], size=16, space_after=14)

card = add_card(s, Inches(0.6), Inches(3.75), Inches(11.9), Inches(2.9))
add_text(s, Inches(0.9), Inches(3.95), Inches(10), Inches(0.4), "VALIDATED AGAINST THE LIVE 30-CAMERA SANDBOX", size=13, color=GREEN, bold=True, font="Consolas")
add_bullets(s, Inches(0.9), Inches(4.45), Inches(11.0), Inches(2.1), [
    "Found and fixed a real cross-thread cv2.VideoCapture segfault during load testing — a genuine stability bug, not a hypothetical one.",
    "Confirmed automatic reconnect/recovery after a forced restart.",
    "Ran a clean, multi-camera, multi-hour unattended session: zero crashes, zero unwanted reconnects, 480+ real detections logged.",
], size=15, space_after=12)
add_footer(s, 5)

# ============================================================ Slide 6: AI analytics — the debugging story
s = add_slide()
add_header(s, "AI-Powered Video Analytics", "Finding — and Fixing — the Plate-Legibility Bottleneck")

steps = [
    "Most sandbox cameras are wide-angle overhead junction cams — genuinely not plate-resolvable at distance. An honest finding, not a failure.",
    "Went back to the camera catalogue's own display names and found a toll-plaza camera (\"Tollnaka\") — camera metadata, not just camera count, was the actionable signal.",
    "On that camera's legible plate crop, our first OCR engine (EasyOCR) capped at ~0.3 confidence after heavy tuning. Swapping to PaddleOCR on the identical crop, no other change: 0.87 confidence.",
    "Even then the real pipeline still failed — root cause was our own bounding-box crop clipping characters. Fixed by padding the detector's box 75% of its own size, verified against real detector output.",
    "Result over a multi-hour unattended run: 480+ detections above threshold across three independently-confirmed cameras — full plausible Indian plates, not fragments (e.g. GJ05AU9828 at 0.96 confidence).",
]
y = Inches(1.55)
for i, step in enumerate(steps):
    num_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.4), Inches(0.4))
    num_circle.fill.solid(); num_circle.fill.fore_color.rgb = ACCENT
    num_circle.line.fill.background(); num_circle.shadow.inherit = False
    tf = num_circle.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1); r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    add_text(s, Inches(1.2), y - Inches(0.02), Inches(11.3), Inches(0.85), step, size=13.5, color=TEXT_MID, line_spacing=1.1)
    y += Inches(1.02)
add_footer(s, 6)

# ============================================================ Slide 7: Detection Feed screenshot
s = add_slide()
add_header(s, "AI Analytics — Live Evidence", "Detection Feed: Watch the AI Working in Real Time")
img_path = EVIDENCE_DIR / "detection_feed.png"
if img_path.exists():
    add_image_framed(s, img_path, Inches(1.1), Inches(1.6), Inches(11.1),
                      caption="Real crop + real annotated frame (bounding box drawn by the detector) for two different live-detected vehicles — not staged.")
add_footer(s, 7)

# ============================================================ Slide 8: Watchlist correlation
s = add_slide()
add_header(s, "Watchlist Correlation", "Real-Time Alerting")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(5.6), Inches(4.5), [
    "Every OCR read — not sampled after the fact — is normalized and matched immediately.",
    "Two-pass matching: exact + OCR-confusion-aware pass (0/O, 1/I, 5/S, 8/B, 2/Z), then fuzzy-similarity fallback.",
    "A match writes an alert with camera, plate, PTS timestamp, confidence, and matched entry — visible on the dashboard within seconds, no manual refresh.",
    "Designed to plug into VAHAN / eGujCop-CCTNS as the real watchlist source in production; a representative watchlist is used for this demo, as explicitly permitted.",
], size=15.5, space_after=16)

img_path = EVIDENCE_DIR / "alert.png"
if img_path.exists():
    add_image_framed(s, img_path, Inches(6.55), Inches(1.6), Inches(6.15),
                      caption="Real watchlist match — 100% confidence, real plate crop, real timestamp.")
add_footer(s, 8)

# ============================================================ Slide 9: Vehicle trace
s = add_slide()
add_header(s, "The Evaluation's Core Test", "Vehicle Route Reconstruction")
add_bullets(s, Inches(0.6), Inches(1.55), Inches(5.6), Inches(2.3), [
    "Given a plate, returns every (camera, location, timestamp) detection in chronological order.",
    "Verified: one real vehicle detected 4 times within 16 seconds, correctly linked across OCR readings of BV2807 and 8V2807 — the confusion-variant matching fix, found via this exact data.",
], size=15, space_after=14)

img_path = EVIDENCE_DIR / "trace.png"
if img_path.exists():
    add_image_framed(s, img_path, Inches(6.55), Inches(1.55), Inches(6.15),
                      caption="Real search result — visual timeline with plate crop.")

card = add_card(s, Inches(0.6), Inches(4.1), Inches(5.6), Inches(2.5), fill=RGBColor(0x1a,0x14,0x08), line_color=AMBER)
add_text(s, Inches(0.85), Inches(4.3), Inches(5.1), Inches(0.35), "HONEST STATUS", size=12, color=AMBER, bold=True, font="Consolas")
add_bullets(s, Inches(0.85), Inches(4.7), Inches(5.1), Inches(1.8), [
    "Three cameras independently confirmed ANPR-viable, ran simultaneously for hours, 480+ real detections.",
    "Comprehensive exact + confusion-variant check: no same-plate sighting across two different cameras yet — real traffic timing and this sandbox's camera mix, not a code limitation.",
], size=13, space_after=8, color=TEXT_MID)
add_footer(s, 9)

# ============================================================ Slide 10: Dashboard
s = add_slide()
add_header(s, "Unified Control-Room View", "The Dashboard")
img_path = EVIDENCE_DIR / "camera_feed.png"
if img_path.exists():
    add_image_framed(s, img_path, Inches(0.7), Inches(1.55), Inches(5.9),
                      caption="Camera grid — live thumbnails, last-seen plate per camera.")
img_path2 = EVIDENCE_DIR / "watchlist.png"
if img_path2.exists():
    add_image_framed(s, img_path2, Inches(6.75), Inches(1.55), Inches(5.9),
                      caption="Watchlist administration.")
add_footer(s, 10)

# ============================================================ Slide 11: Tech stack
s = add_slide()
add_header(s, "Implementation", "Technology Stack")
stack = [
    ("Capture", "OpenCV + FFmpeg backend (Python)"),
    ("Detection", "Ultralytics YOLO — pretrained, CPU-only"),
    ("OCR", "PaddleOCR (PP-OCRv6)"),
    ("Storage", "SQLite (pilot); schema-compatible path to PostgreSQL"),
    ("Dashboard", "Streamlit, custom dark theme"),
]
y = Inches(1.7)
for label, val in stack:
    card = add_card(s, Inches(0.6), y, Inches(11.9), Inches(0.72))
    add_text(s, Inches(0.9), y + Inches(0.13), Inches(2.3), Inches(0.45), label, size=15, color=ACCENT, bold=True, font="Consolas")
    add_text(s, Inches(3.3), y + Inches(0.13), Inches(9), Inches(0.45), val, size=15, color=TEXT_HI)
    y += Inches(0.9)
add_text(s, Inches(0.6), y + Inches(0.1), Inches(11.9), Inches(0.5),
          "No GPU, no external message bus required at this scale — deliberately minimal-dependency for fast, reliable deployment.",
          size=14, color=TEXT_MID)
add_footer(s, 11)

# ============================================================ Slide 12: Scalability
s = add_slide()
add_header(s, "Path to Statewide Scale", "Scalability, Security & Deployment")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.9), Inches(4.8), [
    "Horizontal scale-out: independent capture/ANPR workers, sharded by camera set — no component assumes a single-machine camera count.",
    "Edge pre-filtering + GPU/accelerator path is the realistic route to ~80,000-camera scale; this prototype's frame-sampling and confidence knobs are exactly what moves to edge-side config.",
    "Credentials never logged (verified in code); consume-only design — never pushes to or controls any gateway.",
    "No department's existing VMS, storage, or retention policy is touched.",
    "RBAC, audit logging, and TLS-everywhere are the production hardening path beyond this prototype.",
], size=16, space_after=18)
add_footer(s, 12)

# ============================================================ Slide 13: Impact
s = add_slide()
add_header(s, "Expected Outcomes", "Operational Benefits")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.9), Inches(4.5), [
    "Single pane of glass across departments — without a multi-year federation program.",
    "Automated, continuous watchlist correlation — proactive alerting instead of manual after-the-fact video review.",
    "Fast, low-risk path to statewide expansion: same architecture, more camera shards, no redesign required.",
], size=18, space_after=22)
add_footer(s, 13)

# ============================================================ Slide 14: Roadmap
s = add_slide()
add_header(s, "Honest About Current Scope", "What's Next")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.9), Inches(4.5), [
    "Face Recognition (FRS) and other analytics — explicitly out of scope for this build, bonus territory.",
    "Live VAHAN / eGujCop integration — currently a representative demo watchlist, as explicitly permitted by the challenge rules.",
    "Edge deployment pilot for bandwidth-constrained districts.",
], size=17, space_after=18)
add_footer(s, 14)

# ============================================================ Slide 15: Thank you
s = add_slide()
add_text(s, Inches(0), Inches(3.0), SLIDE_W, Inches(1.0),
          "Thank you", size=48, color=TEXT_HI, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.9), SLIDE_W, Inches(0.5),
          "SENTINEL — Unified Viewing & Metadata Analytics", size=18, color=ACCENT, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved {OUT_PATH} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
